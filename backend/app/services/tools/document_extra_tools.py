"""Sprint D.4: Doküman tool'lari (PDF birleştirme/bölme, PPTX, markdown→html)."""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Any, Dict, List

from app.services.tools.base import BaseTool, ToolContext, ToolResult

logger = logging.getLogger(__name__)


# ============================================================
# pdf_merge
# ============================================================


class PDFMergeTool(BaseTool):
    name = "pdf_merge"
    description = "Birden fazla PDF dosyasini sirayla tek bir PDF'e birlestirir."
    permission = "file_system"
    parameters = {
        "type": "object",
        "properties": {
            "input_paths": {
                "type": "array",
                "items": {"type": "string"},
                "description": "Birlestirilecek PDF dosyalarinin yollari (siraya dikkat)"},
            "output_path": {"type": "string", "description": "Cikti PDF dosyasinin yolu"}},
        "required": ["input_paths", "output_path"]}

    async def execute(self, args: Dict[str, Any], context: ToolContext) -> ToolResult:
        try:
            from pypdf import PdfWriter  # type: ignore
        except ImportError:
            return ToolResult(ok=False, error="pypdf paketi yuklu degil")

        inputs: List[str] = args.get("input_paths") or []
        output = (args.get("output_path") or "").strip()
        if not inputs or not output:
            return ToolResult(ok=False, error="input_paths ve output_path zorunlu")

        writer = PdfWriter()
        added = 0
        try:
            for p in inputs:
                pp = Path(p).expanduser().resolve()
                if not pp.exists():
                    return ToolResult(ok=False, error=f"Bulunamadi: {pp}")
                writer.append(str(pp))
                added += 1
            out_path = Path(output).expanduser().resolve()
            out_path.parent.mkdir(parents=True, exist_ok=True)
            with out_path.open("wb") as f:
                writer.write(f)
        except Exception as exc:
            return ToolResult(ok=False, error=f"PDF birlestirme hata: {exc}")

        size = out_path.stat().st_size
        return ToolResult(
            ok=True,
            output=f"{added} PDF birlestirildi -> {out_path} ({size} bayt)",
            data={"output_path": str(out_path), "merged_count": added, "bytes": size},
        )


# ============================================================
# pdf_split
# ============================================================


class PDFSplitTool(BaseTool):
    name = "pdf_split"
    description = (
        "Bir PDF dosyasini sayfalara veya belirli sayfa araliklarina boler. "
        "ranges parametresi verilmezse her sayfa ayri PDF'e bolunur."
    )
    permission = "file_system"
    parameters = {
        "type": "object",
        "properties": {
            "input_path": {"type": "string"},
            "output_dir": {"type": "string", "description": "Cikti klasoru"},
            "ranges": {
                "type": "array",
                "items": {"type": "string"},
                "description": "(opsiyonel) Sayfa araliklari ['1-3', '5', '7-9']. Bos ise her sayfa ayri."}},
        "required": ["input_path", "output_dir"]}

    async def execute(self, args: Dict[str, Any], context: ToolContext) -> ToolResult:
        try:
            from pypdf import PdfReader, PdfWriter  # type: ignore
        except ImportError:
            return ToolResult(ok=False, error="pypdf paketi yuklu degil")

        in_path = (args.get("input_path") or "").strip()
        out_dir = (args.get("output_dir") or "").strip()
        if not in_path or not out_dir:
            return ToolResult(ok=False, error="input_path ve output_dir zorunlu")

        ip = Path(in_path).expanduser().resolve()
        if not ip.exists():
            return ToolResult(ok=False, error=f"Bulunamadi: {ip}")
        od = Path(out_dir).expanduser().resolve()
        od.mkdir(parents=True, exist_ok=True)

        try:
            reader = PdfReader(str(ip))
            total = len(reader.pages)

            ranges = args.get("ranges") or []
            if not ranges:
                ranges = [str(i + 1) for i in range(total)]

            outputs: List[str] = []
            for idx, r in enumerate(ranges):
                if "-" in r:
                    a, b = r.split("-", 1)
                    start = int(a.strip()) - 1
                    end = int(b.strip())
                else:
                    start = int(r.strip()) - 1
                    end = start + 1
                if start < 0 or end > total or start >= end:
                    return ToolResult(ok=False, error=f"Gecersiz sayfa araligi: {r} (toplam {total})")

                writer = PdfWriter()
                for p in range(start, end):
                    writer.add_page(reader.pages[p])
                out_file = od / f"{ip.stem}_{r.replace('-', '_')}.pdf"
                with out_file.open("wb") as f:
                    writer.write(f)
                outputs.append(str(out_file))
        except Exception as exc:
            return ToolResult(ok=False, error=f"PDF bolme hata: {exc}")

        return ToolResult(
            ok=True,
            output=f"{len(outputs)} parca olusturuldu -> {od}",
            data={"output_dir": str(od), "files": outputs, "total_pages": total},
        )


# ============================================================
# pptx_generate
# ============================================================


class PPTXGenerateTool(BaseTool):
    name = "pptx_generate"
    description = (
        "PowerPoint slayt sunumu (.pptx) olusturur. Basit slayt listesi alir: "
        "[{'title': 'Baslik', 'bullets': ['mad1', 'mad2']}, ...]"
    )
    permission = "file_system"
    parameters = {
        "type": "object",
        "properties": {
            "output_path": {"type": "string"},
            "title": {"type": "string", "description": "Sunum baslik slaytinin baslik metni"},
            "subtitle": {"type": "string", "description": "(opsiyonel) Sunum alt baslik"},
            "slides": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string"},
                        "bullets": {
                            "type": "array",
                            "items": {"type": "string"}}},
                    "required": ["title"]},
                "description": "Icerik slaytlari"}},
        "required": ["output_path", "title", "slides"]}

    async def execute(self, args: Dict[str, Any], context: ToolContext) -> ToolResult:
        try:
            from pptx import Presentation  # type: ignore
        except ImportError:
            return ToolResult(ok=False, error="python-pptx paketi yuklu degil. pip install python-pptx")

        output = (args.get("output_path") or "").strip()
        title = (args.get("title") or "").strip()
        subtitle = (args.get("subtitle") or "").strip()
        slides = args.get("slides") or []
        if not output or not title or not isinstance(slides, list):
            return ToolResult(ok=False, error="output_path, title ve slides zorunlu")

        try:
            prs = Presentation()

            # Title slayti
            tlay = prs.slide_layouts[0]
            tslide = prs.slides.add_slide(tlay)
            tslide.shapes.title.text = title
            if subtitle and len(tslide.placeholders) > 1:
                tslide.placeholders[1].text = subtitle

            # Icerik slaytlari (Title + Content layout)
            content_layout = prs.slide_layouts[1]
            for s in slides:
                if not isinstance(s, dict):
                    continue
                slide = prs.slides.add_slide(content_layout)
                slide.shapes.title.text = str(s.get("title", "Slayt"))
                bullets = s.get("bullets") or []
                if bullets and len(slide.placeholders) > 1:
                    body = slide.placeholders[1].text_frame
                    body.text = str(bullets[0])
                    for b in bullets[1:]:
                        p = body.add_paragraph()
                        p.text = str(b)

            out_path = Path(output).expanduser().resolve()
            out_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(out_path))
            size = out_path.stat().st_size
        except Exception as exc:
            return ToolResult(ok=False, error=f"pptx olusturma hata: {exc}")

        return ToolResult(
            ok=True,
            output=f"PowerPoint olusturuldu: {out_path} ({len(slides) + 1} slayt, {size} bayt)",
            data={"output_path": str(out_path), "slide_count": len(slides) + 1, "bytes": size},
        )


# ============================================================
# markdown_to_html
# ============================================================


class MarkdownToHtmlTool(BaseTool):
    name = "markdown_to_html"
    description = (
        "Markdown metnini GitHub-flavored HTML'e cevirir. Kod blogu syntax highlighting destekler "
        "(markdown + pygments paketi gerekli)."
    )
    permission = "none"
    parameters = {
        "type": "object",
        "properties": {
            "markdown": {"type": "string"},
            "output_path": {
                "type": "string",
                "description": "(opsiyonel) Sonucu bu dosyaya yaz; verilmezse string olarak doner"},
            "wrap_html": {
                "type": "boolean",
                "default": True,
                "description": "<html><body>...</body></html> sarmalayicisi ekle"}},
        "required": ["markdown"]}

    async def execute(self, args: Dict[str, Any], context: ToolContext) -> ToolResult:
        md_text = args.get("markdown") or ""
        if not md_text.strip():
            return ToolResult(ok=False, error="markdown bos olamaz")

        try:
            import markdown  # type: ignore
        except ImportError:
            return ToolResult(ok=False, error="markdown paketi yuklu degil. pip install markdown pygments")

        try:
            html_body = markdown.markdown(
                md_text,
                extensions=[
                    "fenced_code",
                    "tables",
                    "toc",
                    "codehilite"],
                extension_configs={
                    "codehilite": {"guess_lang": True, "css_class": "highlight"}},
            )
        except Exception as exc:
            return ToolResult(ok=False, error=f"markdown parse hata: {exc}")

        wrap = bool(args.get("wrap_html", True))
        if wrap:
            html = (
                "<!DOCTYPE html>\n<html>\n<head>\n"
                "<meta charset='utf-8'>\n"
                "<title>Markdown Output</title>\n"
                "<style>\n"
                "body{font-family:-apple-system,sans-serif;max-width:780px;margin:2rem auto;padding:0 1rem;line-height:1.6;}\n"
                "code{background:#f5f5f5;padding:0.15em 0.3em;border-radius:3px;}\n"
                "pre{background:#272822;color:#f8f8f2;padding:1em;border-radius:6px;overflow-x:auto;}\n"
                "pre code{background:transparent;color:inherit;padding:0;}\n"
                "table{border-collapse:collapse;}\n"
                "th,td{border:1px solid #ddd;padding:0.4em 0.7em;}\n"
                "</style>\n"
                "</head>\n<body>\n"
                + html_body
                + "\n</body>\n</html>"
            )
        else:
            html = html_body

        out_path = (args.get("output_path") or "").strip()
        if out_path:
            try:
                op = Path(out_path).expanduser().resolve()
                op.parent.mkdir(parents=True, exist_ok=True)
                op.write_text(html, encoding="utf-8")
                return ToolResult(
                    ok=True,
                    output=f"HTML yazildi: {op} ({len(html)} char)",
                    data={"output_path": str(op), "bytes": len(html.encode('utf-8'))},
                )
            except Exception as exc:
                return ToolResult(ok=False, error=f"yazma hata: {exc}")

        # Inline donus
        return ToolResult(
            ok=True,
            output=html[:4000] + ("..." if len(html) > 4000 else ""),
            data={"html": html, "length": len(html)},
        )